# --- Standard Libraries ---
import os
import time
import re
import requests
from io import BytesIO

# --- Third-Party Libraries ---
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor
from dotenv import load_dotenv
import google.generativeai as genai
from PIL import Image, ImageDraw, ImageFont

# --- Streamlit Configuration and Gemini Setup ---
def configure_app():
    st.set_page_config(page_title="AI Presentation Architect", page_icon="✨", layout="wide")
    st.title("✨ AI Presentation Architect")
    
    # Load environment variables
    load_dotenv()
    gemini_api_key, pexels_api_key = os.getenv("GEMINI_API_KEY"), os.getenv("PEXELS_API_KEY")
    
    # Validate API keys
    if not gemini_api_key: st.error("GEMINI_API_KEY not found."); st.stop()
    if not pexels_api_key: st.warning("PEXELS_API_KEY not found. Image search disabled.")
    
    # Configure Gemini AI
    genai.configure(api_key=gemini_api_key)
    return pexels_api_key

# --- Image Generation and Fetching ---
def create_placeholder_image(text):
    """Creates a placeholder image with the given text if no real image is found."""
    width, height = 1200, 800
    img = Image.new('RGB', (width, height), color='#E9ECEF')
    draw = ImageDraw.Draw(img)

    # Try to load a readable font
    try: font = ImageFont.truetype("arial.ttf", 60)
    except IOError: font = ImageFont.load_default(size=60)

    # Center align each line of text
    lines = text.split('\n')
    total_h = sum(draw.textbbox((0,0), l, font=font)[3] for l in lines)
    y = (height - total_h) / 2
    for line in lines:
        bbox = draw.textbbox((0, 0), line, font=font, align="center")
        pos = ((width - (bbox[2] - bbox[0])) / 2, y)
        draw.text(pos, line, fill='#6C757D', font=font, align="center")
        y += bbox[3] + 10

    img_buffer = BytesIO()
    img.save(img_buffer, format='PNG')
    img_buffer.seek(0)
    return img_buffer

def fetch_image_from_pexels(api_key, query):
    """Fetches an image from Pexels or falls back to a placeholder."""
    if not api_key or not query:
        return create_placeholder_image(f"Image for:\n{query}")
    headers = {"Authorization": api_key}
    params = {"query": query, "per_page": 1}
    try:
        response = requests.get("https://api.pexels.com/v1/search", headers=headers, params=params, timeout=10)
        response.raise_for_status()
        data = response.json()
        if not data.get("photos"):
            st.warning(f"No image for '{query}'. Using placeholder.")
            return create_placeholder_image(f"No image for:\n'{query}'")
        img_resp = requests.get(data["photos"][0]["src"]["large"], timeout=10)
        img_resp.raise_for_status()
        return BytesIO(img_resp.content)
    except requests.exceptions.RequestException as e:
        st.error(f"Pexels API Error: {e}. Using placeholder.")
        return create_placeholder_image(f"API Error for:\n'{query}'")

# --- AI Text Generation ---
def generate_content_from_ai(prompt):
    """Generates presentation content using Gemini AI."""
    try:
        model = genai.GenerativeModel('gemini-1.5-flash')
        response = model.generate_content(prompt)
        if response.prompt_feedback.block_reason:
            st.error(f"Request blocked by AI safety filters: {response.prompt_feedback.block_reason.name}")
            return None
        return response.text
    except Exception as e:
        st.error(f"AI model error: {e}")
        return None

# --- Parsing AI Response into Structured Slide Data ---
def parse_ai_response(content):
    """Parses the AI output into structured slide dictionaries."""
    slides_data = []
    slide_blocks = content.split('---')
    for block in slide_blocks:
        if not block.strip():
            continue
        slide_dict = {}
        lines = block.strip().split('\n')
        current_key, current_value = "", ""
        for line in lines:
            if ':' in line:
                if current_key:
                    slide_dict[current_key] = current_value.strip()
                key, value = line.split(':', 1)
                current_key, current_value = key.strip().lower().replace(' ', '_'), value.strip()
            else:
                current_value += "\n" + line.strip()
        if current_key:
            slide_dict[current_key] = current_value.strip()
        if 'subtitle' in slide_dict:
            slide_dict['content'] = slide_dict['subtitle']
        if 'title' in slide_dict:
            slides_data.append(slide_dict)
    return slides_data

# --- Utility: Text Wrapping for Long Slide Content ---
def split_text_for_slides(text, max_lines=8, max_chars_per_line=90):
    """Splits long text into multiple chunks for multiple slides."""
    chunks, current_lines = [], []
    text_lines = text.split('\n')
    for line in text_lines:
        while len(line) > max_chars_per_line:
            split_pos = line.rfind(' ', 0, max_chars_per_line)
            split_pos = split_pos if split_pos != -1 else max_chars_per_line
            current_lines.append(line[:split_pos])
            line = line[split_pos:].lstrip()
            if len(current_lines) >= max_lines:
                chunks.append("\n".join(current_lines))
                current_lines = []
        current_lines.append(line)
        if len(current_lines) >= max_lines:
            chunks.append("\n".join(current_lines))
            current_lines = []
    if current_lines:
        chunks.append("\n".join(current_lines))
    return chunks if chunks else [""]

# --- Slide Layout Functions ---
def add_slide_using_template_layout(prs, slide_info, content_chunk, image_stream=None):
    """Adds a slide using existing template layout from .pptx file."""
    layout_names = {
        "Title Slide": ["Title Slide", "Title"],
        "Content Slide": ["Title and Content", "Picture with Caption", "Two Content"],
        "Conclusion Slide": ["Section Header", "Title and Content"]
    }.get(slide_info['slide_type'], ["Title and Content"])

    chosen_layout = next((l for name in layout_names for l in prs.slide_layouts if l.name == name), prs.slide_layouts[1])
    slide = prs.slides.add_slide(chosen_layout)

    # Identify placeholders
    placeholders = {'title': getattr(slide.shapes, 'title', None), 'body': None, 'pic': None}
    for shape in slide.placeholders:
        ph_type = shape.placeholder_format.type
        if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            placeholders['title'] = shape
        elif ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.SUBTITLE):
            placeholders['body'] = shape
        elif ph_type == PP_PLACEHOLDER.PICTURE:
            placeholders['pic'] = shape

    # Insert content
    if placeholders['title']:
        placeholders['title'].text = slide_info['title']
    if placeholders['body']:
        placeholders['body'].text_frame.text = content_chunk
        placeholders['body'].text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        placeholders['body'].text_frame.word_wrap = True
    if image_stream and placeholders['pic']:
        try:
            placeholders['pic'].insert_picture(image_stream)
        except Exception as e:
            st.error(f"Could not insert picture into placeholder: {e}")

def add_title_slide_layout(prs, slide_info, image_stream=None):
    """Adds a custom title slide layout (used when no template is uploaded)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if image_stream:
        try:
            slide.shapes.add_picture(image_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)
        except Exception as e:
            st.error(f"Could not add background image: {e}")

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.75), Inches(9), Inches(1.5))
    tf_title = title_box.text_frame
    tf_title.text = slide_info['title']
    p_title = tf_title.paragraphs[0]
    p_title.font.size, p_title.font.bold, p_title.font.color.rgb, p_title.alignment = Pt(44), True, RGBColor(255, 255, 255), PP_ALIGN.CENTER
    tf_title.word_wrap, tf_title.auto_size = True, MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(9), Inches(1.0))
    tf_sub = subtitle_box.text_frame
    tf_sub.text = slide_info.get('content', '')
    p_sub = tf_sub.paragraphs[0]
    p_sub.font.size, p_sub.font.color.rgb, p_sub.alignment = Pt(24), RGBColor(255, 255, 255), PP_ALIGN.CENTER
    tf_sub.word_wrap, tf_sub.auto_size = True, MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

def add_content_slide_layout(prs, title, content_chunk, image_stream=None):
    """Adds a content slide with optional image and text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1.0))
    p = title_box.text_frame.paragraphs[0]
    p.text, p.font.size, p.font.bold, p.alignment = title, Pt(36), True, PP_ALIGN.LEFT
    title_box.text_frame.word_wrap = True

    content_top_margin = Inches(1.5)
    if image_stream:
        # Dual layout: Text left, image right
        text_container = {"left": Inches(0.5), "top": content_top_margin, "width": Inches(4.5), "height": Inches(5.5)}
        img_container = {"left": Inches(5.5), "top": content_top_margin, "width": Inches(4.0), "height": Inches(5.5)}
        text_box = slide.shapes.add_textbox(**text_container)
        tf = text_box.text_frame
        tf.paragraphs[0].text = content_chunk
        tf.paragraphs[0].font.size = Pt(14)
        tf.word_wrap, tf.auto_size = True, MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        try:
            image_stream.seek(0)
            img = Image.open(image_stream)
            aspect = float(img.width) / img.height
            container_aspect = float(img_container["width"]) / img_container["height"]
            image_stream.seek(0)
            if aspect > container_aspect:
                pic = slide.shapes.add_picture(image_stream, img_container["left"], img_container["top"], width=img_container["width"])
            else:
                pic = slide.shapes.add_picture(image_stream, img_container["left"], img_container["top"], height=img_container["height"])
            pic.left = img_container["left"] + (img_container["width"] - pic.width) / 2
            pic.top = img_container["top"] + (img_container["height"] - pic.height) / 2
        except Exception as e:
            st.error(f"Could not add image: {e}")
    else:
        # Full-width text-only layout
        text_container = {"left": Inches(0.5), "top": content_top_margin, "width": Inches(9), "height": Inches(5.5)}
        text_box = slide.shapes.add_textbox(**text_container)
        tf = text_box.text_frame
        tf.paragraphs[0].text = content_chunk
        tf.paragraphs[0].font.size = Pt(18)
        tf.word_wrap, tf.auto_size = True, MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

# --- Core Logic to Generate PPT from Prompt ---
def generate_ppt(prompt, user_images, auto_image, pexels_key, template_file):
    # Step 1: Generate content using Gemini AI
    with st.spinner("Step 1/3: Generating presentation script with AI..."):
        content = generate_content_from_ai(prompt)
        if not content:
            return None
        slides_data = parse_ai_response(content)
        if not slides_data or len(slides_data) < 5:
            st.error("AI did not return a valid structure.")
            st.info("DEBUG: Raw AI Output:")
            st.code(content if content else "No content.")
            return None

    # Step 2: Create slides in PowerPoint
    prs = Presentation(pptx=template_file) if template_file else Presentation()
    with st.spinner("Step 2/3: Populating slides..."):
        user_image_idx = 0
        for slide_info in slides_data:
            image_to_add = None
            if user_image_idx < len(user_images):
                image_to_add = user_images[user_image_idx]
                user_image_idx += 1
            elif auto_image and pexels_key and slide_info.get('image_search_query', '').lower() not in ['none', 'n/a', 'no image']:
                image_to_add = fetch_image_from_pexels(pexels_key, slide_info['image_search_query'])

            content_chunks = split_text_for_slides(slide_info.get('content', ''))
            for i, chunk in enumerate(content_chunks):
                current_slide_info = slide_info.copy()
                if i > 0:
                    current_slide_info['title'] = f"{slide_info['title']} (cont.)"
                image_for_this_slide = image_to_add if i == 0 else None
                if template_file:
                    add_slide_using_template_layout(prs, current_slide_info, chunk, image_for_this_slide)
                else:
                    if current_slide_info.get('slide_type') == 'Title Slide' and i == 0:
                        add_title_slide_layout(prs, current_slide_info, image_for_this_slide)
                    else:
                        add_content_slide_layout(prs, current_slide_info['title'], chunk, image_for_this_slide)

    # Step 3: Export presentation to buffer
    with st.spinner("Step 3/3: Finalizing presentation..."):
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
    return ppt_buffer

# --- Streamlit User Interface ---
def main():
    pexels_key = configure_app()
    
    # Sidebar settings
    st.sidebar.header("🎨 Design Options")
    template_file = st.sidebar.file_uploader("Upload a Design Template (.pptx)", type=['pptx'])
    st.sidebar.header("🖼️ Image Options")
    auto_image = st.sidebar.checkbox("Find images automatically", value=True)
    user_images = st.sidebar.file_uploader("Or upload your own images", type=['png', 'jpg', 'jpeg'], accept_multiple_files=True)

    # Tabs for topic/text input
    st.header("✍️ Content Creation Mode")
    tab1, tab2 = st.tabs(["💡 Generate from Topic", "📋 Generate from Text"])

    # Prompt templates
    topic_prompt_template = """..."""  # (Left as-is for brevity)
    text_prompt_template = """..."""   # (Left as-is for brevity)

    with tab1:
        topic_input = st.text_input("Enter a topic...")
        if st.button("🚀 Generate from Topic"):
            if topic_input:
                prompt = topic_prompt_template.format(topic=topic_input)
                ppt_buffer = generate_ppt(prompt, user_images, auto_image, pexels_key, template_file)
                if ppt_buffer:
                    st.session_state['generated_ppt'] = ppt_buffer
                    st.session_state['ppt_name'] = f"{topic_input.replace(' ','_')}_Presentation.pptx"
                    st.success("Presentation generated successfully!")
            else:
                st.warning("Please enter a topic.")

    with tab2:
        text_input = st.text_area("Paste your full text...", height=250)
        if st.button("🚀 Generate from Text"):
            if text_input:
                prompt = text_prompt_template.format(text=text_input)
                ppt_buffer = generate_ppt(prompt, user_images, auto_image, pexels_key, template_file)
                if ppt_buffer:
                    st.session_state['generated_ppt'] = ppt_buffer
                    st.session_state['ppt_name'] = "AI_Generated_Presentation.pptx"
                    st.success("Presentation generated successfully!")
            else:
                st.warning("Please paste your text.")

    # Show download button if generation successful
    if 'generated_ppt' in st.session_state and st.session_state['generated_ppt']:
        st.download_button(
            label="⬇️ Download Presentation",
            data=st.session_state['generated_ppt'],
            file_name=st.session_state['ppt_name'],
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            on_click=lambda: st.session_state.update({'generated_ppt': None})
        )

if __name__ == "__main__":
    main()
